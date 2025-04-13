# -*- coding: utf-8 -*-
import csv
import os
from pptx import Presentation
from pypinyin import pinyin, Style
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from gtts import gTTS
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tkinter as tk
from tkinter import filedialog, messagebox
import requests
import threading
from deep_translator import GoogleTranslator

# Function to create a PowerPoint presentation from a template
def create_ppt_from_template(vocab_list, template_path, output_path):
    """
    Create a PowerPoint presentation using a template file.
    """
    print(f"Loading template from: {template_path}")

    if os.path.exists(template_path):
        modification_time = os.path.getmtime(template_path)
        print(f"Template last modified: {modification_time} (Unix timestamp)")

    # Load template and get first slide
    prs = Presentation(template_path)
    template_slide = prs.slides[0]

    # Log placeholders for debugging
    print("Template slide placeholders:")
    for placeholder in template_slide.placeholders:
        print(f"  idx: {placeholder.placeholder_format.idx}, name: {placeholder.name}, type: {placeholder.placeholder_format.type}")

    print("Template slide shapes:")
    for shape in template_slide.shapes:
        if shape.is_placeholder:
            print(f"  Placeholder idx: {shape.placeholder_format.idx}, name: {shape.name}, type: {shape.placeholder_format.type}")
        else:
            print(f"  Shape: {shape.name}, not a placeholder")

    # Loop through each vocabulary pair and create a slide
    for chinese, english in vocab_list:
        # Duplicate the template slide
        slide = prs.slides.add_slide(template_slide.slide_layout)

        # Populate placeholders using their idx values
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:  # Chinese title
                placeholder.text = chinese
            elif placeholder.placeholder_format.idx == 1:  # English subtitle
                placeholder.text = english
            elif placeholder.placeholder_format.idx == 14:  # Pinyin text
                placeholder.text = ' '.join([p[0] for p in pinyin(chinese, style=Style.TONE)])
            elif placeholder.placeholder_format.idx == 15:  # Audio placeholder
                try:
                    # Generate and add TTS audio
                    tts = gTTS(chinese, lang='zh')
                    audio_path = f"media/{chinese}.mp3"
                    tts.save(audio_path)

                    if os.path.exists(audio_path):
                        print(f"Audio file '{audio_path}' created.")
                    else:
                        print(f"Failed to create audio file '{audio_path}'.")

                    # Position audio at placeholder location
                    if placeholder.placeholder_format.idx == 15:
                        left = placeholder.left
                        top = placeholder.top
                        width = placeholder.width
                        height = placeholder.height
                        slide.shapes.add_movie(audio_path, left, top, width=width, height=height)

                    # Add clickable transparent overlay
                    if placeholder.shape_type == MSO_SHAPE_TYPE.MEDIA:
                        left = placeholder.left
                        top = placeholder.top
                        width = placeholder.width
                        height = placeholder.height

                        transparent_shape = slide.shapes.add_shape(
                            MSO_SHAPE_TYPE.RECTANGLE, left, top, width, height
                        )
                        transparent_shape.fill.solid()
                        transparent_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        transparent_shape.fill.transparency = 1.0
                        transparent_shape.line.fill.background()

                except Exception as e:
                    print(f"Error adding audio for '{chinese}': {e}")

        # Add Pixabay images to the left side of the slide
        image_query = english.replace(" ", "+")  # Use the English word as the query
        api_key = "49711697-387fe155204b00a0af7ff7360"  # Replace with your actual API key
        image_url = search_pixabay_images(image_query, api_key)

        if image_url:
            try:
                # Download the image
                image_path = f"media/{image_query}.jpg"
                response = requests.get(image_url, stream=True)
                response.raise_for_status()
                with open(image_path, "wb") as image_file:
                    for chunk in response.iter_content(1024):
                        image_file.write(chunk)

                # Adjust the image position and size to be larger and more centered, but slightly to the left
                left = Inches(1.5)  # Slightly more to the left
                top = Inches(1.5)   # Centered vertically
                slide.shapes.add_picture(image_path, left, top, width=Inches(5), height=Inches(5))

                # Remove the attribution text below the image
                # Removed the 'Image by Pixabay' text addition logic

                print(f"Added image for '{english}' to the slide.")
            except Exception as e:
                print(f"Error adding image for '{english}': {e}")
        else:
            # Use the placeholder image if no image is found
            placeholder_path = "media/placeholder-image.png"
            if os.path.exists(placeholder_path):
                left = Inches(1.5)  # Slightly more to the left
                top = Inches(1.5)   # Centered vertically
                slide.shapes.add_picture(placeholder_path, left, top, width=Inches(5), height=Inches(5))
                print(f"Added placeholder image for '{english}' to the slide.")
            else:
                print(f"Placeholder image not found at {placeholder_path}. Skipping image addition.")

    # Remove the original template slide safely
    xml_slides = prs.slides._sldIdLst  # Access the slide ID list
    slides = list(xml_slides)  # Convert to a list for iteration
    xml_slides.remove(slides[0])  # Remove the first slide (template slide)

    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

    # Auto-delete the generated mp3 files
    for file_name in os.listdir("media"):
        if file_name.endswith(".mp3"):
            file_path = os.path.join("media", file_name)
            try:
                os.remove(file_path)
                print(f"Deleted audio file: {file_path}")
            except Exception as e:
                print(f"Error deleting file {file_path}: {e}")

    # Auto-delete the downloaded images after saving the presentation, except for placeholder-image.png
    for file_name in os.listdir("media"):
        if file_name.endswith(".jpg") and file_name != "placeholder-image.png":
            file_path = os.path.join("media", file_name)
            try:
                os.remove(file_path)
                print(f"Deleted image file: {file_path}")
            except Exception as e:
                print(f"Error deleting file {file_path}: {e}")

# Function to search for images using the Pixabay API
def search_pixabay_images(query, api_key):
    base_url = "https://pixabay.com/api/"
    # Properly encode the query to handle spaces and special characters
    encoded_query = query.replace(" ", "+")
    params = {
        "key": api_key,
        "q": encoded_query,  # Use the encoded query
        "image_type": "illustration",  # Fetch illustration images
        "safesearch": "true"
    }

    try:
        response = requests.get(base_url, params=params)
        response.raise_for_status()  # Raise an error for HTTP issues
        data = response.json()

        if "hits" in data and len(data["hits"]) > 0:
            return data["hits"][0]["largeImageURL"]  # Return the first image URL
        else:
            print(f"No images found for query: {query}")
            return None
    except requests.exceptions.RequestException as e:
        print(f"Error fetching images from Pixabay: {e}")
        return None

# Function to translate Chinese text to English
def translate_chinese_to_english(text):
    """Translate Chinese to English using Google's translation service."""
    try:
        translator = GoogleTranslator(source='zh-CN', target='en')
        translation = translator.translate(text)
        return translation
    except Exception as e:
        print(f"Translation error: {e}")
        return f"[Translation failed: {text}]"

def run_gui():
    def select_csv():
        file_path = filedialog.askopenfilename(filetypes=[("CSV Files", "*.csv")])
        if file_path:
            csv_path_var.set(file_path)

    # Set the default template path
    default_template_path = "template.pptx"

    def select_template():
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            template_path_var.set(file_path)
        else:
            # If no custom template is selected, use the default template
            template_path_var.set(default_template_path)

    # Create the Tk root window first
    root = tk.Tk()
    root.title("Mandarin Vocabulary Presentation Generator")
    root.geometry("750x500")  # Increased height from 450 to 500 to give more vertical space
    root.configure(bg="#f4f4f9")  # Set a light corporate background color
    root.resizable(False, False)  # Disable resizing to make the window fixed size

    # Pre-fill the template path with the default template
    template_path_var = tk.StringVar(value=default_template_path)
    
    # Define use_default_template before it's used
    use_default_template = tk.BooleanVar(value=True)

    # Ensure language_var is defined before toggle_template is called
    language_var = tk.StringVar(value="en")

    # Update the toggle_template function to use language_var safely
    def toggle_template():
        if use_default_template.get():
            template_path_var.set(default_template_path)
            template_label.config(
                text=("PowerPoint 模板文件 (默认)" if language_var.get() == "zh" else "PowerPoint Template File (Default)"),
                fg="#28a745"
            )  # Green for default
            template_entry.config(state="disabled")
            browse_button.config(state="disabled")
        else:
            template_label.config(
                text=("PowerPoint 模板文件 (自定义)" if language_var.get() == "zh" else "PowerPoint Template File (Custom)"),
                fg="#333333"
            )  # Normal color for custom
            template_entry.config(state="normal")
            browse_button.config(state="normal")

    # Highlight the default template in the GUI
    def update_template_highlight():
        if template_path_var.get() == default_template_path:
            template_label.config(fg="#333333", font=("Helvetica", 12))  # Normal text style for default
        else:
            template_label.config(fg="#333333", font=("Helvetica", 12))  # Normal style for custom
        
        # Remove the call to update_language to avoid circular dependency

    # Update the highlight whenever the template path changes
    template_path_var.trace_add("write", lambda *args: update_template_highlight())

    # Define template_label before calling update_template_highlight
    template_label = tk.Label(
        root,
        text="Select PowerPoint Template:",
        font=("Helvetica", 12),
        fg="#333333",
        bg="#f4f4f9"
    )
    template_label.grid(row=4, column=0, sticky="e", padx=10, pady=5)

    # Call the highlight function initially to set the default state
    update_template_highlight()

    def generate_ppt_thread():
        def task():
            # Disable the Generate button and remove hover effects to prevent multiple generations
            generate_button.config(state="disabled")
            generate_button.unbind("<Enter>")
            generate_button.unbind("<Leave>")

            csv_path = csv_path_var.get()
            template_path = template_path_var.get()
            output_path = output_path_var.get()

            if not csv_path or not template_path or not output_path:
                messagebox.showerror("Error", "Please fill in all fields.")
                generate_button.config(state="normal")  # Re-enable the button
                generate_button.bind("<Enter>", on_enter_green)
                generate_button.bind("<Leave>", on_leave_green)
                return

            if not output_path.lower().endswith(".pptx"):
                output_path += ".pptx"

            try:
                # Show a loading message
                loading_label = tk.Label(
                    root,
                    text="Generating PowerPoint... Please wait.",
                    font=("Helvetica", 12),
                    fg="#333333",
                    bg="#f4f4f9"
                )
                loading_label.grid(row=9, column=0, columnspan=3, pady=10)
                root.update_idletasks()

                vocab = []
                with open(csv_path, mode="r", encoding="utf-8") as file:
                    reader = csv.reader(file)
                    for row in reader:
                        # Handle different CSV formats based on auto-translate setting
                        if len(row) >= 2:
                            # Traditional format with Chinese and English columns
                            chinese, english = row[0].strip(), row[1].strip()
                            vocab.append((chinese, english))
                        elif len(row) == 1 and auto_translate_var.get():
                            # Chinese-only format with auto-translation enabled
                            chinese = row[0].strip()
                            # Update loading message to show translation is happening
                            loading_label.config(text=f"Translating: {chinese}...")
                            root.update_idletasks()
                            # Translate Chinese to English
                            english = translate_chinese_to_english(chinese)
                            print(f"Translated '{chinese}' to '{english}'")
                            vocab.append((chinese, english))
                        else:
                            # Skip invalid rows
                            print(f"Skipping invalid row: {row}")
                            continue

                # Update loading message to show we're generating the PPT
                loading_label.config(text="Generating PowerPoint... Please wait.")
                root.update_idletasks()

                create_ppt_from_template(vocab, template_path, output_path)

                # Remove the loading message
                loading_label.destroy()

                # Open the presentation automatically
                try:
                    os.startfile(output_path)
                    print(f"Opening PowerPoint file: {output_path}")
                except Exception as e:
                    print(f"Error opening PowerPoint file: {e}")

                messagebox.showinfo("Success", f"Presentation saved to {output_path}")
            except Exception as e:
                # Remove the loading message in case of an error
                loading_label.destroy()
                messagebox.showerror("Error", f"An error occurred: {e}")
            finally:
                # Re-enable the Generate button and restore hover effects
                generate_button.config(state="normal")
                generate_button.bind("<Enter>", on_enter_green)
                generate_button.bind("<Leave>", on_leave_green)

        # Run the task in a separate thread
        threading.Thread(target=task).start()

    # Adjust title label
    title_label = tk.Label(
        root,
        text="Mandarin Vocabulary Presentation Generator",
        font=("Helvetica", 16, "bold"),
        fg="#333333",
        bg="#f4f4f9"
    )
    title_label.grid(row=0, column=0, columnspan=3, pady=10)

    # Adjust separator
    separator = tk.Frame(root, height=2, bd=0, relief="solid", bg="#cccccc")
    separator.grid(row=1, column=0, columnspan=3, sticky="we", pady=10)

    # Adjust CSV file input
    csv_label = tk.Label(
        root,
        text="Vocabulary CSV File:",
        font=("Helvetica", 12),
        fg="#333333",
        bg="#f4f4f9"
    )
    csv_label.grid(row=2, column=0, sticky="e", padx=10, pady=5)
    csv_path_var = tk.StringVar()
    tk.Entry(
        root,
        textvariable=csv_path_var,
        width=50,
        font=("Helvetica", 10),
        bd=2,
        relief="groove"
    ).grid(row=2, column=1, pady=5)
    csv_browse_button = tk.Button(
        root,
        text="Browse",
        command=select_csv,
        font=("Helvetica", 10),
        bg="#0078d7",
        fg="white",
        relief="flat",
        padx=10
    )
    csv_browse_button.grid(row=2, column=2, padx=10, pady=5)

    # Add auto-translation checkbox
    auto_translate_var = tk.BooleanVar(value=False)
    auto_translate_checkbox = tk.Checkbutton(
        root,
        text="Auto-translate Chinese to English",
        variable=auto_translate_var,
        font=("Helvetica", 10),
        bg="#f4f4f9"
    )
    auto_translate_checkbox.grid(row=3, column=1, pady=5, sticky="w")

    # Adjust template file input
    template_entry = tk.Entry(
        root,
        textvariable=template_path_var,
        width=50,
        font=("Helvetica", 10),
        bd=2,
        relief="groove"
    )
    template_entry.grid(row=4, column=1, pady=5)

    browse_button = tk.Button(
        root,
        text="Browse",
        command=select_template,
        font=("Helvetica", 10),
        bg="#0078d7",
        fg="white",
        relief="flat",
        padx=10
    )
    browse_button.grid(row=4, column=2, padx=10, pady=5)

    # Define the checkbox for toggling the default template
    checkbox = tk.Checkbutton(
        root,
        text="Use Default Template (Recommended)",
        variable=use_default_template,
        command=toggle_template,
        font=("Helvetica", 10),
        bg="#f4f4f9"
    )
    checkbox.grid(row=5, column=1, pady=5, sticky="w")  # Now in its own row (5)

    # Adjust output file name input
    output_label = tk.Label(
        root,
        text="Output File Name:",
        font=("Helvetica", 12),
        fg="#333333",
        bg="#f4f4f9"
    )
    output_label.grid(row=6, column=0, sticky="e", padx=10, pady=5)  
    output_path_var = tk.StringVar()
    tk.Entry(
        root,
        textvariable=output_path_var,
        width=50,
        font=("Helvetica", 10),
        bd=2,
        relief="groove"
    ).grid(row=6, column=1, pady=5)  

    # Adjust generate button
    generate_button = tk.Button(
        root,
        text="Generate Presentation",
        command=generate_ppt_thread,
        font=("Helvetica", 12, "bold"),
        bg="#28a745",
        fg="white",
        relief="flat",
        padx=15,
        pady=5
    )
    generate_button.grid(row=7, column=0, columnspan=3, pady=20)  

    # Update instructions to reflect auto-translation feature
    instructions_label = tk.Label(
        root,
        text=(
            "CSV Format: Each row should contain either:\n"
            "1. Two columns: Chinese and English words, separated by a comma (e.g., 你好,Hello)\n"
            "2. One column with only Chinese words if auto-translation is enabled (e.g., 你好)\n"
            "\nExample:\n"
            "你好,Hello\n"
            "谢谢,Thank you\n"
            "再见,Goodbye\n"
        ),
        font=("Helvetica", 10),
        fg="#333333",
        bg="#f4f4f9",
        wraplength=700,
        justify="left",
    )
    instructions_label.grid(row=8, column=0, columnspan=3, pady=10)

    # Adjust footer
    footer_label = tk.Label(
        root,
        text="Developed for Professional Use",
        font=("Helvetica", 10, "italic"),
        fg="#666666",
        bg="#f4f4f9"
    )
    footer_label.grid(row=9, column=0, columnspan=3, pady=10)  # Changed from row 8 to row 9

    # Add a language toggle button
    def toggle_language():
        if language_var.get() == "en":
            language_var.set("zh")
            update_language("zh")
        else:
            language_var.set("en")
            update_language("en")

    # Ensure update_language is defined before update_template_highlight
    def update_language(lang):
        if lang == "zh":
            title_label.config(text="普通话词汇演示文稿生成器")
            csv_label.config(text="词汇 CSV 文件：")
            template_label.config(text="选择 PowerPoint 模板：")  # This is the key line for translation
            checkbox.config(text="使用默认模板（推荐）")
            browse_button.config(text="浏览")
            csv_browse_button.config(text="浏览")
            output_label.config(text="输出文件名：")
            generate_button.config(text="生成演示文稿")
            auto_translate_checkbox.config(text="自动翻译中文为英文")  # Added translation for auto-translate checkbox
            instructions_label.config(
                text=(
                    "CSV 格式：每行应包含以下内容之一：\n"
                    "1. 两列：中文和英文单词，用逗号分隔（如：你好,Hello）\n"
                    "2. 如果启用自动翻译，则只需一列中文单词（如：你好）\n"
                    "\n示例：\n"
                    "你好,Hello\n"
                    "谢谢,Thank you\n"
                    "再见,Goodbye\n"
                )
            )
            footer_label.config(text="为专业用途开发")
            
            # Force update the template label based on current state
            if use_default_template.get():
                template_label.config(text="选择 PowerPoint 模板：", fg="#28a745")
            else:
                template_label.config(text="选择 PowerPoint 模板：", fg="#333333")
                
        else:
            title_label.config(text="Mandarin Vocabulary Presentation Generator")
            csv_label.config(text="Vocabulary CSV File:")
            template_label.config(text="Select PowerPoint Template:")  # This is the key line for translation
            checkbox.config(text="Use Default Template (Recommended)")
            browse_button.config(text="Browse")
            csv_browse_button.config(text="Browse")
            output_label.config(text="Output File Name:")
            generate_button.config(text="Generate Presentation")
            auto_translate_checkbox.config(text="Auto-translate Chinese to English")  # Added translation for auto-translate checkbox
            instructions_label.config(
                text=(
                    "CSV Format: Each row should contain either:\n"
                    "1. Two columns: Chinese and English words, separated by a comma (e.g., 你好,Hello)\n"
                    "2. One column with only Chinese words if auto-translation is enabled (e.g., 你好)\n"
                    "\nExample:\n"
                    "你好,Hello\n"
                    "谢谢,Thank you\n"
                    "再见,Goodbye\n"
                )
            )
            footer_label.config(text="Developed for Professional Use")
            
            # Force update the template label based on current state
            if use_default_template.get():
                template_label.config(text="Select PowerPoint Template:", fg="#28a745")
            else:
                template_label.config(text="Select PowerPoint Template:", fg="#333333")

    language_var = tk.StringVar(value="en")
    language_button = tk.Button(
        root,
        text="中文/English",
        command=toggle_language,
        font=("Helvetica", 10),
        bg="#0078d7",
        fg="white",
        relief="flat",
        padx=10
    )
    language_button.grid(row=9, column=2, pady=10, sticky="e")

    # Add hover effect for buttons
    def on_enter(e):
        e.widget["bg"] = "#0056a3"  # Darker blue for hover

    def on_leave(e):
        e.widget["bg"] = "#0078d7"  # Original blue

    def on_enter_green(e):
        e.widget["bg"] = "#1e7e34"  # Darker green for hover

    def on_leave_green(e):
        e.widget["bg"] = "#28a745"  # Original green

    # Apply hover effects to buttons
    for button in [
        language_button,
        generate_button,
        csv_browse_button,  # CSV Browse button
        browse_button   # Template Browse button
    ]:
        if button == generate_button:
            button.bind("<Enter>", on_enter_green)
            button.bind("<Leave>", on_leave_green)
        else:
            button.bind("<Enter>", on_enter)
            button.bind("<Leave>", on_leave)

    # Configure grid weights for proper responsiveness
    for i in range(9):  # Adjust row weights
        root.grid_rowconfigure(i, weight=1)

    for i in range(3):  # Adjust column weights
        root.grid_columnconfigure(i, weight=1)

    # Initialize the toggle state
    toggle_template()
    
    # Ensure template entry and browse button are initially disabled when default is selected
    if use_default_template.get():
        template_entry.config(state="disabled")
        browse_button.config(state="disabled")

    root.mainloop()

if __name__ == "__main__":
    run_gui()